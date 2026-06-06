"""
Shared fixtures for the DPA test suite.

Import order matters: env vars must be set BEFORE importing backend modules
because db_connector.py and auth_service.py call sys.exit() at module load
if required vars are missing.
"""
import os
import sys
import pathlib
import tempfile
import shutil

# ── 1. Set env vars before any backend import ──────────────────────────────
TEST_ROOT = pathlib.Path(__file__).parent
FIXTURES  = TEST_ROOT / "fixtures"
FIXTURES.mkdir(exist_ok=True)
(FIXTURES / "output").mkdir(exist_ok=True)
(FIXTURES / "images").mkdir(exist_ok=True)
(FIXTURES / "doc").mkdir(exist_ok=True)

os.environ.setdefault("DB_HOST",          "localhost")
os.environ.setdefault("DB_PORT",          "5432")
os.environ.setdefault("DB_NAME",          "DPA_TEST")
os.environ.setdefault("DB_USER",          "postgres")
os.environ.setdefault("DB_PASSWORD",      "postgres_test")
os.environ.setdefault("JWT_SECRET_KEY",   "test-jwt-secret-key-minimum-32-chars-long")
os.environ.setdefault("TEMPLATE_PATH",    str(FIXTURES / "test_template.pptx"))
os.environ.setdefault("OUTPUT_DIR",       str(FIXTURES / "output"))
os.environ.setdefault("IMAGE_WIN_ROOT",   r"D:\Auto_detect\Result")
os.environ.setdefault("IMAGE_MOUNT_ROOT", str(FIXTURES / "images"))
os.environ.setdefault("DOC_ROOT",         str(FIXTURES / "doc"))
os.environ.setdefault("ALLOWED_ORIGINS",  "http://localhost:5190")
os.environ.setdefault("COOKIE_SECURE",    "false")

# ── 2. Add backend/ to sys.path ─────────────────────────────────────────────
BACKEND = pathlib.Path(__file__).parent.parent / "backend"
sys.path.insert(0, str(BACKEND))

import pytest
from unittest.mock import MagicMock, patch
from starlette.testclient import TestClient

# ── 3. Lazy app import (after env vars are set) ──────────────────────────────
@pytest.fixture(scope="session")
def app():
    from main import app as _app
    if hasattr(_app, "state") and hasattr(_app.state, "limiter"):
        _app.state.limiter.enabled = False
    from routers.auth import limiter as auth_limiter
    auth_limiter.enabled = False
    return _app


@pytest.fixture(scope="session")
def client(app):
    with TestClient(app, raise_server_exceptions=True) as c:
        yield c


# ── 4. JWT helpers ────────────────────────────────────────────────────────────
@pytest.fixture(scope="session")
def make_token():
    from services.auth_service import create_access_token
    def _make(user_id="EMP001", name="Test User", role="QA Engineer"):
        return create_access_token({"sub": user_id, "name": name, "role": role})
    return _make


@pytest.fixture
def admin_token(make_token):
    return make_token(user_id="admin", name="Admin User", role="admin")


@pytest.fixture
def qa_token(make_token):
    return make_token(user_id="EMP001", name="QA Engineer", role="QA Engineer")


@pytest.fixture
def user_token(make_token):
    return make_token(user_id="EMP002", name="Regular User", role="user")


@pytest.fixture
def auth_headers(qa_token):
    return {"Authorization": f"Bearer {qa_token}"}


@pytest.fixture
def admin_headers(admin_token):
    return {"Authorization": f"Bearer {admin_token}"}


# ── 5. Cookie-based auth (mirrors real frontend behaviour) ───────────────────
@pytest.fixture
def auth_cookies(qa_token):
    return {"dpa_token": qa_token}


@pytest.fixture
def admin_cookies(admin_token):
    return {"dpa_token": admin_token}


# ── 6. DB mock helpers ────────────────────────────────────────────────────────
def _make_cursor(rows=None, fetchone_val=None):
    """Build a psycopg2-style mock cursor."""
    cur = MagicMock()
    cur.__enter__ = lambda s: s
    cur.__exit__ = MagicMock(return_value=False)
    cur.fetchall.return_value = rows or []
    cur.fetchone.return_value = fetchone_val
    return cur


def _make_conn(cursor=None):
    """Build a psycopg2-style mock connection."""
    conn = MagicMock()
    conn.cursor.return_value = cursor or _make_cursor()
    conn.__enter__ = lambda s: s
    conn.__exit__ = MagicMock(return_value=False)
    return conn


@pytest.fixture
def mock_db(monkeypatch):
    """Patch DBConnector so no real DB is needed. Returns (mock_conn, mock_cursor)."""
    from services import db_connector
    cursor = _make_cursor()
    conn   = _make_conn(cursor)
    monkeypatch.setattr(db_connector.DBConnector, "get_dpa_connection",
                        staticmethod(lambda: conn))
    monkeypatch.setattr(db_connector.DBConnector, "release_dpa_connection",
                        staticmethod(lambda c: None))
    return conn, cursor


# ── 7. Standard test-data fixtures ───────────────────────────────────────────
SAMPLE_USER_ROW = {
    "user_id":       "EMP001",
    "full_name":     "QA Engineer One",
    "role":          "QA Engineer",
    "password_hash": "$2b$12$KIXjV3qJ8Z1mN2pL5oR7OuWvHkT9xYcBdEfGhAiJsPlQrMnUzSw4K",
}

SAMPLE_PR_ROW = {
    "pr_no":                    "PR2024001",
    "excel_file_name":          "PR2024001.xlsx",
    "request_date":             "2024-01-15",
    "subject":                  "DPA Report for MT0 QFN32 Package",
    "purpose":                  "Reliability evaluation at T0",
    "conclusion":               "",
    "customer_name":            "MT0",
    "assembly_site":            "Hana Microelectronics",
    "package_type":             "QFN",
    "date_code":                "2401",
    "package_size":             "5x5 mm",
    "number_of_lot":            "1",
    "pin_ball_count":           "32",
    "requestor_name_dept":      "QA Department",
    "reliability_staff_name":   "Reliability Staff A",
    "rel_request_number_1":     "REL-2024-001",
    "rel_request_number_2":     "",
    "rel_request_number_3":     "",
    "order_lot":                "MTDQS0906.1",
    "cust_assy":                "MT0-QFN32-5x5",
    "device":                   "MT001",
    "die_size":                 "3.0x3.0 mm",
    "dap_size":                 "3.1x3.1 mm",
    "lf_stock_no":              "LF-001",
    "die_attach_material":      "Ag Paste",
    "wire_type":                "Au 25um",
    "mold_compound":            "EME-7351",
    "plating_finish":           "NiPdAu",
}

SAMPLE_HISTORY_ROW = {
    "id":         1,
    "pr_no":      "PR2024001",
    "order_lot":  "MTDQS0906.1",
    "revision":   "A",
    "timepoint":  "T0",
    "user_id":    "EMP001",
    "file_name":  "DPA_Report_PR2024001_T0_MTDQS0906.1_20240115_103045.pptx",
    "file_path":  str(FIXTURES / "output" / "DPA_Report_PR2024001_T0_MTDQS0906.1_20240115_103045.pptx"),
    "created_at": "2024-01-15T10:30:45",
}


@pytest.fixture
def sample_user():
    return dict(SAMPLE_USER_ROW)


@pytest.fixture
def sample_pr():
    return dict(SAMPLE_PR_ROW)


@pytest.fixture
def sample_history():
    return dict(SAMPLE_HISTORY_ROW)


# ── 8. Temp output dir (cleaned after each test) ─────────────────────────────
@pytest.fixture
def tmp_output(tmp_path, monkeypatch):
    monkeypatch.setenv("OUTPUT_DIR", str(tmp_path))
    import services.report_generator as rg
    monkeypatch.setattr(rg, "OUTPUT_DIR", str(tmp_path))
    return tmp_path


# ── 9. Minimal PPTX template fixture ─────────────────────────────────────────
@pytest.fixture(scope="session")
def minimal_template(tmp_path_factory):
    """
    Create a minimal .pptx with one slide containing key placeholders.
    Used by PPTX generation tests so the real template file is not required.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])

    tf = slide.shapes.add_textbox(Inches(0), Inches(0), Inches(4), Inches(1)).text_frame
    tf.text = "{background_info.customer_name}"

    tf2 = slide.shapes.add_textbox(Inches(0), Inches(1), Inches(4), Inches(1)).text_frame
    tf2.text = "{bill_of_materials.order_lot}"

    tf3 = slide.shapes.add_textbox(Inches(0), Inches(2), Inches(4), Inches(1)).text_frame
    tf3.text = "{Image_records.EXTERNAL_1-1}"

    out = tmp_path_factory.mktemp("template") / "test_template.pptx"
    prs.save(str(out))
    return out

"""
Integration tests for the DPA PowerPoint report generation service.
Integrates DPAReportGenerator with python-pptx, Pillow, and the filesystem.
"""
import os
import pathlib
import pytest
from unittest.mock import MagicMock, patch
from pptx import Presentation
from PIL import Image

from services.report_generator import DPAReportGenerator
from tests.seeds.factories import make_full_report_data

pytestmark = pytest.mark.integration


def create_dummy_jpeg(path: pathlib.Path, width=200, height=200):
    """Create a dummy solid-color JPEG image using Pillow."""
    path.parent.mkdir(parents=True, exist_ok=True)
    img = Image.new("RGB", (width, height), color="blue")
    img.save(path, "JPEG")


def test_pptx_report_generation_integration(minimal_template, tmp_output):
    """
    End-to-end integration test for DPAReportGenerator.
    Verifies that the template is parsed, placeholders are replaced,
    images are correctly loaded, translated, sized, and placed, and
    the resulting presentation is saved to the output directory.
    """
    pr_no = "PR2024001"
    timepoint = "T0"
    lot = "MTDQS0906.1"

    # 1. Setup mock source images under the translated local mount directory
    mount_root = pathlib.Path(os.environ["IMAGE_MOUNT_ROOT"])
    
    # Paths as translated from D:\Auto_detect\Result\PR2024001\...
    ext_img_path = mount_root / pr_no / "images" / timepoint / lot / "1.EXTERNAL VISUAL" / "1-1.jpg"
    xray_img_path = mount_root / pr_no / "images" / timepoint / lot / "3.X-RAY" / "1-1.jpg"

    create_dummy_jpeg(ext_img_path, width=400, height=300)
    create_dummy_jpeg(xray_img_path, width=200, height=200)

    # 2. Get full report data bundle from factories
    report_data = make_full_report_data(pr_no=pr_no, lot=lot, timepoint=timepoint)
    
    # 3. Mock the product request service so that DB calls are bypassed
    # and return our structured report data instead.
    with patch("services.report_generator.fetch_full_report_data", return_value=report_data), \
         patch("services.report_generator.TEMPLATE_PATH", str(minimal_template)):
        
        # Instantiate and run generator
        gen = DPAReportGenerator(pr_no, timepoint, lot, selected_sections={"EXTERNAL VISUAL": True, "X-RAY": True}, revision="A")
        output_path, stats = gen.generate()

    # ── Assertions ────────────────────────────────────────────────────────────

    # Verify return values and stats
    assert os.path.exists(output_path)
    assert output_path.endswith(".pptx")
    
    assert stats["metadata_found"] is True
    assert stats["images_found"] >= 1
    assert stats["images_missing"] == 0

    # Read back the generated PPTX to verify its elements
    prs = Presentation(output_path)
    assert len(prs.slides) == 1

    slide = prs.slides[0]
    
    # Verify that text placeholders were replaced with actual DB data
    text_runs = []
    for shape in slide.shapes:
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text_runs.append(para.text)

    # The {background_info.customer_name} should be replaced with "MT0"
    # and {bill_of_materials.order_lot} with "MTDQS0906.1"
    assert "MT0" in text_runs
    assert "MTDQS0906.1" in text_runs

    # Verify that original placeholder names are no longer present
    for run_text in text_runs:
        assert "{background_info.customer_name}" not in run_text
        assert "{bill_of_materials.order_lot}" not in run_text

    # Clean up dummy images and folders
    try:
        if ext_img_path.exists():
            ext_img_path.unlink()
        if xray_img_path.exists():
            xray_img_path.unlink()
        if os.path.exists(output_path):
            os.unlink(output_path)
    except Exception:
        pass

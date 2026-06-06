import os
import re
import pathlib
from datetime import datetime
from functools import lru_cache

from pptx import Presentation
from pptx.util import Inches
from pptx.oxml.ns import qn
from PIL import Image as PILImage

from services.product_request_service import fetch_full_report_data
from logger import get_logger

log = get_logger("report_generator")

TEMPLATE_PATH = os.environ.get("TEMPLATE_PATH", r"D:\DPA\Template\DPA report.pptx")
OUTPUT_DIR    = os.environ.get("OUTPUT_DIR",    r"D:\DPA\output")

_IMAGE_WIN_ROOT   = os.getenv("IMAGE_WIN_ROOT",   r"D:\Auto_detect\Result")
_IMAGE_MOUNT_ROOT = os.getenv("IMAGE_MOUNT_ROOT", _IMAGE_WIN_ROOT)


def _translate_image_path(path: str | None) -> str | None:
    """Translate a DB-stored Windows path to the actual path in this environment."""
    if not path:
        return path
    norm_path = path.replace("\\", "/")
    norm_win  = _IMAGE_WIN_ROOT.replace("\\", "/")
    if norm_path.lower().startswith(norm_win.lower()):
        relative = norm_path[len(norm_win):].lstrip("/")
        return str(pathlib.PurePosixPath(_IMAGE_MOUNT_ROOT) / relative)
    return path


@lru_cache(maxsize=512)
def _get_image_dimensions(img_path: str) -> tuple[int, int]:
    """Return (width, height) in pixels; cached to avoid repeated PIL opens."""
    try:
        with PILImage.open(img_path) as img:
            return img.size
    except Exception as e:
        log.warning("Cannot read image dimensions for %s: %s", img_path, e)
        return (1, 1)

# ─────────────────────────────────────────────────────────────────────────────
# Image size config per category key.
# format: "KEY": (width_in, height_in, margin_in)
#   None → auto-detect from Container Rectangle in template
# ─────────────────────────────────────────────────────────────────────────────
IMAGE_SIZE_CONFIG = {
    "EXTERNAL":  None,
    "EXTERANAL": None,   # legacy typo in some PPTX templates
    "XRAY":      (0.82, 0.82, 0.02),
    "DELAM":     (2.80, 1.30, 0.02),
    "DECAP":     (0.98, 0.80, 0.02),
    "IMC":       (0.90, 0.90, 0.02),
    "C-R":       (0.90, 0.90, 0.02),
    "CR":        None,
    "IMAGE":     (1.50, 2.00, 0.02),
    "UNIT":      (0.90, 0.90, 0.02),
    "BS":        None,
}


class DPAReportGenerator:
    def __init__(self, pr_number, timepoint, lot, selected_sections, revision="A"):
        self.pr_number = pr_number
        self.timepoint = timepoint
        self.lot = lot
        self.selected_sections = selected_sections
        self.revision = revision
        self.db_data = None
        self.prs = None
        self.unique_units = []
        self.current_sem_page = 0
        self.stats = {
            "metadata_found": False,
            "images_found":   0,
            "images_missing": 0,
            "total_slides":   0,
            "missing_list":   [],
        }

    def load_data(self):
        log.info("Loading data  PR=%s  TP=%s  Lot=%s", self.pr_number, self.timepoint, self.lot)
        self.db_data = fetch_full_report_data(self.pr_number, self.lot, self.timepoint)

        meta = self.db_data.get("metadata", {})
        if meta:
            self.stats["metadata_found"] = True
            log.info("Metadata loaded: pr_no=%s  customer=%s  lot=%s",
                     meta.get("pr_no"), meta.get("customer_name"), meta.get("order_lot"))
        else:
            log.warning("No metadata found for PR=%s — text placeholders will be empty", self.pr_number)

        images = self.db_data.get("images", {})
        log.info("Images in DB: %d entries  keys=%s", len(images), list(images.keys())[:10])

        imc = self.db_data.get("imc", [])
        log.info("IMC records: %d", len(imc))

        sem = self.db_data.get("sem_records", [])
        log.info("SEM records: %d", len(sem))

        if not os.path.exists(TEMPLATE_PATH):
            log.error("Template not found: %s", TEMPLATE_PATH)
            raise FileNotFoundError(f"Template not found at {TEMPLATE_PATH}")
        self.prs = Presentation(TEMPLATE_PATH)
        self.stats["total_slides"] = len(self.prs.slides)
        log.info("Template loaded: %d slides", self.stats["total_slides"])

    def generate(self):
        self.load_data()

        # Pre-calculate categories for all slides before processing
        self._slide_categories = {}
        for s in self.prs.slides:
            self._slide_categories[s.slide_id] = self._identify_slide_category(s)

        # Extract unique unit_ids and sort them
        sem_list = self.db_data.get("sem_records", []) if self.db_data else []
        unit_ids = list(set(str(item.get("unit_id", "")) for item in sem_list if item.get("unit_id")))
        def get_unit_num(uid):
            num_part = "".join(filter(str.isdigit, uid))
            return int(num_part) if num_part else 9999
        self.unique_units = sorted(unit_ids, key=get_unit_num)
        log.info("CROSS SECTION unique units: %s", self.unique_units)

        slides_to_delete = []
        for i, slide in enumerate(self.prs.slides):
            category = self._slide_categories.get(slide.slide_id)
            if category and not self.selected_sections.get(category, False):
                slides_to_delete.append(slide)
                log.debug("Slide %d (%s) will be deleted at the end", i, category)

        # Dynamic slide duplication for CROSS SECTION INSPECTION (Done BEFORE deleting to avoid part naming conflicts)
        if self.selected_sections.get("CROSS SECTION INSPECTION", False):
            sem_slide_idx = -1
            for i, slide in enumerate(self.prs.slides):
                if self._slide_categories.get(slide.slide_id) == "CROSS SECTION INSPECTION":
                    sem_slide_idx = i
                    break
            
            if sem_slide_idx != -1:
                num_sem_slides_needed = max(1, (len(self.unique_units) + 1) // 2)
                log.info("Need %d CROSS SECTION slides for %d units", num_sem_slides_needed, len(self.unique_units))
                for _ in range(num_sem_slides_needed - 1):
                    self._duplicate_slide(sem_slide_idx)

        # Process all slides except the ones to be deleted
        for slide in self.prs.slides:
            if any(slide is s for s in slides_to_delete):
                continue
            self._process_slide(slide)

        # Defer unselected slides deletion to the very end to avoid part naming collision
        if slides_to_delete:
            log.info("Deleting %d unselected slides at the end", len(slides_to_delete))
            for slide_to_del in slides_to_delete:
                for idx, slide in enumerate(self.prs.slides):
                    if slide is slide_to_del:
                        self._delete_slide(idx)
                        break

        os.makedirs(OUTPUT_DIR, exist_ok=True)

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        filename = f"DPA_Report_{self.pr_number}_{self.timepoint}_{self.lot}_{timestamp}.pptx"
        output_path = os.path.join(OUTPUT_DIR, filename)
        self.prs.save(output_path)

        log.info("Report saved: %s  (images found=%d  missing=%d)",
                 filename, self.stats["images_found"], self.stats["images_missing"])
        if self.stats["missing_list"]:
            log.warning("Missing image placeholders: %s", self.stats["missing_list"])
        return output_path, self.stats

    def _delete_slide(self, index):
        """Remove a slide fully from the presentation (ID list + XML part)."""
        xml_slides = self.prs.slides._sldIdLst
        slide_id_el = xml_slides[index]
        rId = slide_id_el.get(qn('r:id'))

        xml_slides.remove(slide_id_el)

        # Remove the actual slide XML part so the file doesn't bloat
        prs_part = self.prs.part
        if hasattr(prs_part, 'drop_rel'):
            prs_part.drop_rel(rId)
        elif rId in prs_part._rels:
            del prs_part._rels[rId]

    def _duplicate_slide(self, source_slide_idx):
        """Duplicate a slide at the specified index and insert it right after."""
        import copy
        source_slide = self.prs.slides[source_slide_idx]
        
        # Create a new blank slide with the same layout
        slide_layout = source_slide.slide_layout
        dest_slide = self.prs.slides.add_slide(slide_layout)
        
        # Delete any default shapes that the layout added automatically
        for shape in list(dest_slide.shapes):
            dest_slide.shapes._spTree.remove(shape._element)
            
        # Copy all shapes from source_slide to dest_slide
        for shape in source_slide.shapes:
            new_el = copy.deepcopy(shape._element)
            dest_slide.shapes._spTree.append(new_el)
            
        # Move the new slide from the end of the presentation to right after the source slide
        xml_slides = self.prs.slides._sldIdLst
        slide_id_el = xml_slides[-1]
        xml_slides.remove(slide_id_el)
        xml_slides.insert(source_slide_idx + 1, slide_id_el)

        # Store category for the duplicated slide
        source_cat = self._slide_categories.get(source_slide.slide_id)
        self._slide_categories[dest_slide.slide_id] = source_cat

    def _identify_slide_category(self, slide):
        """Identify which DB category this slide belongs to by checking its placeholders."""
        for shape in slide.shapes:
            if not shape.has_text_frame: continue
            txt = shape.text.strip()
            
            # Check Image_records
            if "{Image_records." in txt:
                clean_key = re.sub(r'^\{Image_records\.\s*', '', txt, flags=re.IGNORECASE)
                clean_key = re.sub(r'\}$', '', clean_key).strip()
                cat_abbr = clean_key.split("_")[0].strip().upper()
                return self._CATEGORY_MAP.get(cat_abbr)
            
            # Check sem_records (should identify as CROSS SECTION INSPECTION)
            if "{sem_records." in txt:
                return "CROSS SECTION INSPECTION"
        return None

    def _process_slide(self, slide):
        category = self._slide_categories.get(slide.slide_id)
        if category == "CROSS SECTION INSPECTION":
            sem_slides = []
            for s in self.prs.slides:
                if self._slide_categories.get(s.slide_id) == "CROSS SECTION INSPECTION":
                    sem_slides.append(s)
            
            page_idx = 0
            for idx, s in enumerate(sem_slides):
                if s is slide:
                    page_idx = idx
                    break
            self.current_sem_page = page_idx
            log.info("Processing CROSS SECTION slide page %d", self.current_sem_page)

        for shape in list(slide.shapes):
            if shape.has_text_frame:
                if "{Image_records." in shape.text:
                    self._insert_image_to_shape(slide, shape)
                elif "{sem_records." in shape.text:
                    # Treat as image only if it's a point-specific record (Unit + Point)
                    # and the shape looks like a standalone image placeholder.
                    # Remove all whitespace/newlines for matching
                    txt_clean = "".join(shape.text.split()).lower()
                    
                    # If it's a dedicated image slot (contains only the placeholder)
                    if txt_clean.startswith("{sem_records.") and txt_clean.endswith("}"):
                        # Extract the inner part
                        inner = txt_clean.strip("{}").lower()
                        col = None
                        for c in ["magnification", "accel_volt", "unit_id", "point_id", "file_path", "image"]:
                            if f"sem_records.{c}" in inner:
                                col = c
                                break
                        
                        if col:
                            remaining = inner.replace(f"sem_records.{col}", "").strip("_ ")
                            r_parts = [p for p in remaining.split("_") if p]
                            if len(r_parts) >= 2: # Has Unit and Point -> It's an image
                                self._insert_image_to_shape(slide, shape)
                                continue

                    self._replace_text_in_shape(shape)
                else:
                    self._replace_text_in_shape(shape)
            if shape.has_table:
                self._replace_text_in_table(shape.table)

    # ------------------------------------------------------------------
    # Text replacement helpers
    # ------------------------------------------------------------------

    def _replace_in_paragraph(self, para):
        """
        Replace placeholders in one paragraph, correctly handling the case
        where a placeholder is split across multiple runs by PowerPoint.
        Merges all run text → applies mapping → writes back into first run.
        """
        if not para.runs:
            return
        full_text = "".join(run.text for run in para.runs)
        if "{" not in full_text or "}" not in full_text:
            return
        new_text = self._map_placeholder_to_value(full_text)
        if new_text == full_text:
            return
        para.runs[0].text = new_text
        for run in para.runs[1:]:
            run.text = ""

    def _replace_text_in_shape(self, shape):
        for paragraph in shape.text_frame.paragraphs:
            self._replace_in_paragraph(paragraph)

    def _replace_text_in_table(self, table):
        """
        Replace placeholders inside table cells while preserving formatting.
        Walks every <a:tc> element directly to handle merged-cell layouts.
        """
        A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
        tbl_el = table._tbl
        seen_tc_ids = set()

        for tc in tbl_el.iter(f"{{{A_NS}}}tc"):
            tc_id = id(tc)
            if tc_id in seen_tc_ids:
                continue
            seen_tc_ids.add(tc_id)

            for p_el in tc.findall(f".//{{{A_NS}}}p"):
                runs = p_el.findall(f"{{{A_NS}}}r")
                if not runs:
                    continue
                full_text = "".join((r.findtext(f"{{{A_NS}}}t") or "") for r in runs)
                if "{" not in full_text or "}" not in full_text:
                    continue
                new_text = self._map_placeholder_to_value(full_text)
                if new_text == full_text:
                    continue
                first_t = runs[0].find(f"{{{A_NS}}}t")
                if first_t is not None:
                    first_t.text = new_text
                for run in runs[1:]:
                    t_el = run.find(f"{{{A_NS}}}t")
                    if t_el is not None:
                        t_el.text = ""

    def _map_placeholder_to_value(self, text: str) -> str:
        """Map every {table.column} placeholder to its DB value (case-insensitive)."""
        if "{" not in text or "}" not in text:
            return text

        meta = self.db_data.get("metadata", {})
        placeholders = re.findall(r"\{[^\}]+\}", text)

        for p in placeholders:
            p_clean = "".join(p.strip("{}").split()).lower()
            
            # --- 1. Special Placeholder: excel_file_name (strip extension) ---
            if "excel_file_name" in p_clean:
                val = meta.get("excel_file_name", "")
                if val:
                    if isinstance(val, str):
                        for ext in [".xlsx", ".xlsm", ".xls"]:
                            if val.lower().endswith(ext):
                                val = val[:-len(ext)]
                                break
                    text = text.replace(p, str(val))
                else:
                    text = text.replace(p, "")
                continue
            
            if p_clean.startswith("values_records."):
                # Remove prefix and any internal spaces (like in "values_records. IMC_1_1")
                val_key = p_clean.replace("values_records.", "").replace(" ", "").upper()
                # Handle IMC values (mapping unit_id like '1-1' to placeholder 'IMC_1_1')
                if val_key.startswith("IMC_"):
                    imc_list = self.db_data.get("imc", [])
                    imc_found = False
                    for item in imc_list:
                        # Normalize unit_id from DB (e.g. '1-1') to match placeholder format (e.g. '1_1')
                        db_unit_id = str(item.get('unit_id', '')).replace('-', '_')
                        if f"IMC_{db_unit_id}" == val_key:
                            val = item.get('imc_percent', "")
                            # Format value as % if numeric
                            val_str = f"{val}%" if (val is not None and val != "") else ""
                            text = text.replace(p, val_str)
                            imc_found = True
                            break
                    if not imc_found:
                        text = text.replace(p, "")
                else:
                    text = text.replace(p, "")
                continue

            # --- 2. Handle sem_records (Magnification, Voltage for SEM) ---
            if p_clean.startswith("sem_records."):
                sem_list = self.db_data.get("sem_records", [])
                
                # Identify column: magnification, accel_volt, point_id, unit_id, file_path, image
                col = None
                for c in ["magnification", "accel_volt", "point_id", "unit_id", "file_path", "image"]:
                    if c in p_clean:
                        col = c
                        break
                
                if col:
                    # Extract Unit and Point from e.g. {sem_records.magnification_1_B1}
                    # Remaining part after 'sem_records.column_'
                    rem_p = p_clean.replace(f"sem_records.{col}", "").strip("_ ")
                    r_parts = [rp.strip() for rp in rem_p.split("_") if rp.strip()]
                    target_unit = r_parts[0] if len(r_parts) > 0 else None
                    target_point = r_parts[1].upper() if len(r_parts) > 1 else None
                    
                    target_idx = -1
                    if target_unit:
                        try:
                            local_idx = int("".join(filter(str.isdigit, str(target_unit)))) - 1
                            target_idx = self.current_sem_page * 2 + local_idx
                        except ValueError:
                            pass
                    
                    target_unit_id = None
                    if 0 <= target_idx < len(self.unique_units):
                        target_unit_id = self.unique_units[target_idx]
                    
                    found = False
                    if target_unit_id:
                        for item in sem_list:
                            db_unit_raw = str(item.get("unit_id", "")).upper()
                            if db_unit_raw == target_unit_id.upper():
                                if target_point:
                                    db_point = str(item.get("point_id", "")).upper()
                                    if db_point == target_point or db_point.replace("-", "") == target_point.replace("-", ""):
                                        val = item.get(col, "")
                                        if col == "magnification" and val:
                                            val = f"{val}X"
                                        elif col == "accel_volt" and val:
                                            val = f"{val}kV"
                                        text = text.replace(p, str(val))
                                        found = True
                                        break
                                else:
                                    val = item.get(col, "")
                                    text = text.replace(p, str(val))
                                    found = True
                                    break
                    
                    if not found:
                        text = text.replace(p, "")
                continue

            # กรณีพิเศษ: Timepoint และ Revision (รวมถึงกรณีสะกดผิดใน Template)
            if "timepoint" in p_clean or "time_point" in p_clean:
                text = text.replace(p, str(self.timepoint))
                continue
            if "revision" in p_clean:
                text = text.replace(p, str(self.revision))
                continue
            if "current_monthyear" in p_clean or "currentmonthyear" in p_clean:
                current_time = datetime.now().strftime("%B %Y")
                text = text.replace(p, current_time)
                continue

            # Try to find match in metadata (e.g., pr_no, bwip.package_code)
            found = False

            # Special case for request_date formatting (MMDD'YY)
            if "request_date" in p_clean and meta.get("request_date"):
                try:
                    rd = meta["request_date"]
                    if isinstance(rd, str):
                        # Try to parse if it's a string
                        dt = datetime.strptime(rd, "%Y-%m-%d")
                    else:
                        dt = rd
                    formatted_date = dt.strftime("%b%d'%y")
                    text = text.replace(p, formatted_date)
                    found = True
                except Exception:
                    pass

            # Prioritize exact match if not already handled
            if not found and p_clean in meta:
                val = meta[p_clean]
                text = text.replace(p, str(val) if val is not None else "")
                found = True
            
            if not found:
                # Fallback to partial match (legacy support)
                for key, val in meta.items():
                    col_name = key.lower().replace(" ", "")
                    # Also strip spaces from placeholder for comparison
                    p_no_space = p_clean.replace(" ", "")
                    if col_name in p_no_space:
                        text = text.replace(p, str(val) if val is not None else "")
                        found = True
                        break
            
            if not found:
                log.warning("Placeholder not matched: %s", p)
                text = text.replace(p, "")

        return text

    # ------------------------------------------------------------------
    # Image insertion
    # ------------------------------------------------------------------

    def _get_container_rect(self, slide, shape):
        """Find the smallest background Rectangle spatially containing the placeholder."""
        cx = shape.left + shape.width  // 2
        cy = shape.top  + shape.height // 2
        best, best_area = None, float("inf")

        for s in slide.shapes:
            if s is shape or s.has_table:
                continue
            if s.has_text_frame and s.text_frame.text.strip():
                continue
            if (s.left <= cx <= s.left + s.width and s.top <= cy <= s.top + s.height):
                area = s.width * s.height
                if area < best_area:
                    best_area = area
                    best = s

        if best:
            w_ratio    = best.width  / max(shape.width,  1)
            h_ratio    = best.height / max(shape.height, 1)
            area_ratio = (shape.width * shape.height) / max(best.width * best.height, 1)
            if w_ratio <= 4.0 and h_ratio <= 2.5 and area_ratio >= 0.12:
                return best.left, best.top, best.width, best.height

        return shape.left, shape.top, shape.width, shape.height

    def _fit_image_in_slot(self, img_path, slot_l, slot_t, slot_w, slot_h, margin=27000):
        """Compute (left, top, width, height) maintaining native aspect ratio."""
        iw, ih = _get_image_dimensions(img_path)

        avail_w = max(slot_w - 2 * margin, 1)
        avail_h = max(slot_h - 2 * margin, 1)

        if iw > 0 and ih > 0:
            img_ar   = iw / ih
            avail_ar = avail_w / avail_h
            if img_ar > avail_ar:
                final_w = avail_w
                final_h = int(round(avail_w / img_ar))
            else:
                final_h = avail_h
                final_w = int(round(avail_h * img_ar))
        else:
            final_w, final_h = avail_w, avail_h

        final_l = slot_l + (slot_w - final_w) // 2
        final_t = slot_t + (slot_h - final_h) // 2
        return final_l, final_t, final_w, final_h

    def _insert_image_to_shape(self, slide, shape):
        placeholder_text = shape.text.strip()
        is_sem_placeholder = bool(re.match(r'^\{sem_records\.', placeholder_text, re.IGNORECASE))

        # Remove prefix and suffix, and also remove any spaces after the dot
        clean_key = re.sub(r'^\{(Image_records|sem_records)\.\s*', '', placeholder_text, flags=re.IGNORECASE)
        clean_key = re.sub(r'\}$', '', clean_key).strip()
        # Remove any internal spaces (like in "IMC _2_1")
        clean_key = clean_key.replace(" ", "")

        image_path = _translate_image_path(self._find_image_path(clean_key))

        if not (image_path and os.path.exists(image_path)):
            self.stats["images_missing"] += 1
            self.stats["missing_list"].append(clean_key)
            log.debug("Image not found: key=%s  resolved_path=%s", clean_key, image_path)
            shape.text = ""
            return

        parts      = clean_key.upper().split("_")
        cat_prefix = parts[0].strip()
        if is_sem_placeholder:
            # Enlarge SEM images to a premium 2.7" x 2.025" with optimized spacing
            override = (2.7, 2.025, 0.0)
            
            lower_txt = placeholder_text.lower()
            if "b1-" in lower_txt:
                box_cx = int(Inches(5.50))
            elif "b2-" in lower_txt:
                box_cx = int(Inches(8.50))
            elif "b3-" in lower_txt:
                box_cx = int(Inches(11.50))
            else:
                box_cx = shape.left + shape.width // 2

            if "image_1_" in lower_txt:
                box_cy = int(Inches(2.84))
            elif "image_2_" in lower_txt:
                box_cy = int(Inches(5.45))
            else:
                box_cy = shape.top + shape.height // 2
        else:
            override = IMAGE_SIZE_CONFIG.get(cat_prefix, None)
            box_cx = shape.left + shape.width  // 2
            box_cy = shape.top  + shape.height // 2

        if override is not None:
            ow_in, oh_in = override[0], override[1]
            om_in        = override[2] if len(override) > 2 else 0.03
            ow_emu = int(Inches(ow_in))
            oh_emu = int(Inches(oh_in))
            om_emu = int(Inches(om_in))
            slot_l = box_cx - ow_emu // 2
            slot_t = box_cy - oh_emu // 2
            slot_w, slot_h = ow_emu, oh_emu
            margin_emu = om_emu
        else:
            c_l, c_t, c_w, c_h = self._get_container_rect(slide, shape)
            shape_cx     = shape.left + shape.width  // 2
            container_cx = c_l + c_w // 2
            PAIR_GAP  = 18000
            TOLERANCE = 30000

            diff = shape_cx - container_cx
            if diff < -TOLERANCE:
                slot_w = (c_w - PAIR_GAP) // 2
                slot_l, slot_t, slot_h = c_l, c_t, c_h
            elif diff > TOLERANCE:
                slot_w = (c_w - PAIR_GAP) // 2
                slot_l = c_l + slot_w + PAIR_GAP
                slot_t, slot_h = c_t, c_h
            else:
                slot_l, slot_t, slot_w, slot_h = c_l, c_t, c_w, c_h
            margin_emu = 27000

        final_l, final_t, final_w, final_h = self._fit_image_in_slot(
            image_path, slot_l, slot_t, slot_w, slot_h, margin=margin_emu
        )
        shape.text = ""
        slide.shapes.add_picture(image_path, final_l, final_t, final_w, final_h)
        self.stats["images_found"] += 1
        log.debug("Image inserted: key=%s  file=%s", clean_key, os.path.basename(image_path))

    # ------------------------------------------------------------------
    # Image path resolution
    # ------------------------------------------------------------------

    _CATEGORY_MAP = {
        "EXTERNAL":  "EXTERNAL VISUAL",
        "EXTERANAL": "EXTERNAL VISUAL",   # legacy typo in some PPTX templates
        "VISUAL":    "EXTERNAL VISUAL",
        "DELAM":     "DELAM",
        "XRAY":      "X-RAY",
        "DECAP":     "DECAP",
        "IMC":       "IMC",
        "CR":        "C-R",
        "C-R":       "C-R",
        "BS":        "BS",    # prefix — matches both BS,WP,SP and BS,WS,SP
        "IMAGE":     "BS",    # bond test images accessed via image_ prefix
        "UNIT":      "BS",    # bond test unit images
    }

    def _find_image_path(self, key: str):
        """Map a template placeholder key to an image file path."""
        images = self.db_data.get("images", {})
        key_upper = key.strip().upper()

        for img_key, path in images.items():
            if key_upper in img_key:
                return path

        # --- Handle sem_records. ---
        # Note: key might be stripped of prefix by _insert_image_to_shape
        lower_key = key.lower()
        is_sem = lower_key.startswith("sem_records.") or "_" in lower_key # Heuristic if already stripped
        
        if is_sem:
            sem_list = self.db_data.get("sem_records", [])
            # Handle both "sem_records.unit_id_1_B1-1" and "unit_id_1_B1-1"
            p_clean = lower_key.replace("sem_records.", "").strip("_ ")
            parts = p_clean.split("_")
            
            # Identify the field/column (e.g. unit_id, magnification)
            # and extract target unit/point
            col = None
            for c in ["magnification", "accel_volt", "unit_id", "point_id", "file_path", "image"]:
                if c in p_clean:
                    col = c
                    break
            
            if col:
                # Remaining after field name
                rem = p_clean.replace(col, "").strip("_ ")
                r_parts = rem.split("_")
                
                if len(r_parts) >= 2:
                    target_unit = r_parts[0]
                    target_point = r_parts[1].upper()
                    
                    target_idx = -1
                    if target_unit:
                        try:
                            local_idx = int("".join(filter(str.isdigit, str(target_unit)))) - 1
                            target_idx = self.current_sem_page * 2 + local_idx
                        except ValueError:
                            pass
                    
                    target_unit_id = None
                    if 0 <= target_idx < len(self.unique_units):
                        target_unit_id = self.unique_units[target_idx]
                    
                    if target_unit_id:
                        for item in sem_list:
                            db_unit_raw = str(item.get("unit_id", "")).upper()
                            if db_unit_raw == target_unit_id.upper():
                                db_point = str(item.get("point_id", "")).upper()
                                if db_point == target_point or db_point.replace("-", "") == target_point.replace("-", ""):
                                    return item.get("file_path")
        
        if lower_key.startswith("sem_records."):
            return None # If it was a sem_records key but no image found, don't fallback to image_records
            
        parts    = [p.strip() for p in key_upper.split("_")]
        cat_abbr = parts[0]
        db_cat   = self._CATEGORY_MAP.get(cat_abbr, "").upper()
        if not db_cat:
            return None

        candidates = []

        if len(parts) == 2:
            try:
                candidates.append(str(int(parts[1])))
            except ValueError:
                candidates.append(parts[1])

        elif len(parts) == 3:
            try:
                unit = str(int(parts[1]))
                pos  = str(int(parts[2]))
            except ValueError:
                unit, pos = parts[1], parts[2]
            candidates.append(f"{unit}-{pos}")
            if cat_abbr == "DECAP":
                try:
                    candidates.append(str((int(unit) - 1) * 5 + int(pos)))
                except ValueError:
                    pass

        for db_seq in candidates:
            for img_key, path in images.items():
                img_cat, sep, img_seq = img_key.partition("_")
                if sep and img_cat.upper().startswith(db_cat) and img_seq == db_seq:
                    return path

        return None
